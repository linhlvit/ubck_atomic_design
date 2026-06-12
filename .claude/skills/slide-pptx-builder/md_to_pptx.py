#!/usr/bin/env python3
"""
md_to_pptx.py — Convert Markdown slide content to PowerPoint (.pptx)
using the VIB/FSS brand template backgrounds.

All positions/sizes measured directly from VIB original PPTX shapes.

Usage:
    python3 md_to_pptx.py --input slides.md --output output.pptx
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# ── Asset paths ────────────────────────────────────────────────────────────────

SKILL_DIR  = Path(__file__).parent
ASSETS_DIR = SKILL_DIR / 'assets'
BG_COVER   = ASSETS_DIR / 'bg_cover.jpeg'    # gradient green→purple, hexagons, FSS logo, www.fss.com.vn
BG_CONTENT = ASSETS_DIR / 'bg_content.jpeg'  # white, border frame, hexagons top-right, FSS logo, www.fss.com.vn

# ── Brand colors (from VIB FSS palette) ───────────────────────────────────────

C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_DARK   = RGBColor(0x26, 0x26, 0x26)
C_TITLE_DARK  = RGBColor(0x40, 0x40, 0x40)
C_GREEN       = RGBColor(0x77, 0xBB, 0x43)   # FSS green
C_PURPLE      = RGBColor(0x5B, 0x45, 0x9E)   # FSS purple
C_PLACEHOLDER = RGBColor(0xCC, 0xCC, 0xCC)

# ── Slide dimensions (from VIB PPTX: 18288000 x 10287000 EMU) ─────────────────

SLIDE_W = Emu(18288000)
SLIDE_H = Emu(10287000)

def i(inches):
    """Convert inches to EMU."""
    return Emu(round(inches * 914400))


# ── Exact VIB measurements (inches, from shape XML) ───────────────────────────
#
# COVER  : title x=0.35" y=1.70" w=13.5" h=3.75"  sz=64pt  (VIB: x=-0.451 y=1.698 sz=8000)
# SECTION: title x=1.42" y=3.79" w=17.33" h=1.31" sz=64pt  (VIB: sz=8000)
# CONTENT: num   x=0.40" y=0.22" w=1.14"  h=0.91" sz=38pt  (VIB: sz=4800)
#          title x=1.66" y=0.22" w=18.87" h=0.91" sz=38pt
# TABLE  : same as CONTENT header
# CLOSING: title x=1.28" y=4.30" w=17.23" h=1.01" sz=38pt  (VIB: sz=4800)
# TOC    : heading x=2.75" y=0.51" w=11.52" h=0.64" sz=26pt (VIB: sz=3200)
#          underline y=0.79" h=0.05"
#          items y=1.87", 2.59", 3.29", 4.08", 4.86" (step ~0.72")
#          numbers x=3.10" y=item_y+0.02" w=0.42" h=0.42" sz=19pt (VIB: sz=2400)


# ── Helper: add full-slide background image ────────────────────────────────────

def add_bg(slide, img_path):
    """Insert image as bottom-layer background covering the entire slide."""
    pic = slide.shapes.add_picture(str(img_path), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    sp_tree = slide.shapes._spTree
    el = pic._element
    sp_tree.remove(el)
    sp_tree.insert(2, el)   # insert after grpSpPr, before all other shapes
    return pic


# ── Helper: add text box ───────────────────────────────────────────────────────

def textbox(slide, x, y, w, h, text,
            size=Pt(18), bold=False, color=C_TEXT_DARK,
            align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in str(text).split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text       = line
        run.font.size  = size
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name  = 'Calibri'
        first = False
    return tb


# ── Helper: solid rectangle ────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


# ── Helper: image placeholder box ─────────────────────────────────────────────

def placeholder(slide, x, y, w, h, desc):
    box = rect(slide, x, y, w, h, fill=C_PLACEHOLDER)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {desc} ]"
    run.font.size   = Pt(14)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    run.font.name   = 'Calibri'
    return box


# ── Helper: parse "1.1 | Title text" ──────────────────────────────────────────

def split_num_title(title):
    m = re.match(r'^([\d.]+)\s*\|\s*(.+)$', title)
    return (m.group(1).strip(), m.group(2).strip()) if m else ('', title)


# ── Helper: draw content header (number badge + title) ────────────────────────

def content_header(slide, title_full):
    """Draw the section-number + title row used by content/table slides.

    VIB exact positions (from slide 33 / slide 5 XML):
      Number: x=0.40" y=0.22" w=1.14" h=0.91" sz=38pt
      Title:  x=1.66" y=0.22" w=18.87" h=0.91" sz=38pt
    The FSS logo in bg_content.jpeg occupies top-left ~(0,0)-(1.5",0.6").
    The number box starts at x=0.40" which partly overlaps the logo area —
    this is intentional in the VIB original (logo is small and number is bold).
    """
    num, title = split_num_title(title_full)
    if num:
        textbox(slide, i(0.40), i(0.22), i(1.14), i(0.91),
                num, size=Pt(38), bold=True, color=C_TITLE_DARK)
        tx = i(1.66)
        tw = SLIDE_W - tx - i(0.30)
    else:
        tx, tw = i(1.66), SLIDE_W - i(1.66) - i(0.30)
    textbox(slide, tx, i(0.22), tw, i(0.91),
            title, size=Pt(38), bold=True, color=C_TITLE_DARK)
    # Green separator line — VIB: just below title row
    rect(slide, i(0.40), i(1.13), SLIDE_W - i(0.80), Emu(22000), fill=C_GREEN)


# ── Markdown parser ────────────────────────────────────────────────────────────

def parse_md(text):
    lines = text.splitlines()
    meta, slides_raw = {}, []
    current, i_line = None, 0

    if lines and lines[0].strip() == '---':
        i_line = 1
        while i_line < len(lines) and lines[i_line].strip() != '---':
            m = re.match(r'^(\w+):\s*(.+)$', lines[i_line])
            if m:
                meta[m.group(1)] = m.group(2).strip()
            i_line += 1
        i_line += 1

    for line in lines[i_line:]:
        m = re.match(r'^##\s+\[(\w[\w-]*)\]\s*$', line)
        if m:
            if current is not None:
                slides_raw.append(current)
            current = {'layout': m.group(1).lower(), 'lines': []}
        elif current is not None:
            current['lines'].append(line)
    if current is not None:
        slides_raw.append(current)
    return meta, slides_raw


def parse_content(raw):
    layout = raw['layout']
    lines  = [l for l in raw['lines'] if l.strip()]
    c = {'layout': layout}

    if layout == 'cover':
        c['title'] = next((l[2:].strip() for l in lines if l.startswith('# ')), '')

    elif layout == 'toc':
        c['items'] = [l.lstrip('- ').strip() for l in lines if l.strip().startswith('-')]

    elif layout == 'section':
        c['title'] = next((l[2:].strip() for l in lines if l.startswith('# ')), '')

    elif layout in ('content', 'content-text'):
        title, body_lines, img = '', [], None
        for l in lines:
            if l.startswith('### '):
                title = l[4:].strip()
            elif re.match(r'^\[image:', l) and l.endswith(']'):
                img = l[7:-1].strip()
            else:
                body_lines.append(l)
        c['title'] = title
        c['body']  = '\n'.join(body_lines).strip()
        c['image'] = img

    elif layout == 'table':
        title, header, rows, got_header = '', [], [], False
        for l in lines:
            if l.startswith('### '):
                title = l[4:].strip()
            elif l.startswith('|'):
                cells = [x.strip() for x in l.strip('|').split('|')]
                if all(re.match(r'^[-:]+$', x) for x in cells if x):
                    continue
                if not got_header:
                    header = cells; got_header = True
                else:
                    rows.append(cells)
        c['title'] = title; c['header'] = header; c['rows'] = rows

    elif layout == 'closing':
        c['title'] = '\n'.join(lines).strip()

    return c


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_cover(prs, c, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COVER)

    # Client label — sits just above the main title
    client = meta.get('client', '')
    if client:
        textbox(slide, i(0.35), i(1.25), i(13), i(0.40),
                client, size=Pt(18), color=C_WHITE)

    # Main title — VIB exact: y=1.698" sz=64pt
    textbox(slide, i(0.35), i(1.70), i(13.5), i(3.75),
            c.get('title', ''), size=Pt(64), bold=True, color=C_WHITE)
    return slide


def build_toc(prs, c, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_CONTENT)

    # Heading — VIB: x=2.75" y=0.51" sz=26pt
    textbox(slide, i(2.75), i(0.51), i(11.52), i(0.64),
            'NỘI DUNG CHÍNH', size=Pt(26), bold=True, color=C_TITLE_DARK)

    # Green underline — VIB: y=0.79" h=0.05"
    rect(slide, i(1.08), i(0.79), i(3.79), i(0.05), fill=C_GREEN)

    # Items — VIB y positions: 1.87, 2.59, 3.29, 4.08, 4.86 (step ~0.72")
    vib_y = [1.87, 2.59, 3.29, 4.08, 4.86]
    items = c.get('items', [])

    for idx, item in enumerate(items):
        if idx >= 6:
            break
        # Calculate y: use VIB positions for first 5, extend by 0.72" steps after
        if idx < len(vib_y):
            item_y = vib_y[idx]
        else:
            item_y = vib_y[-1] + (idx - len(vib_y) + 1) * 0.72

        # Number badge — left side, then item text to the right
        # VIB original has ovals at x=3.10" but those are overlapping/animated;
        # reconstructed: badge at x=2.09", item text starts at x=2.70"
        badge = rect(slide, i(2.09), i(item_y + 0.04), i(0.42), i(0.42), fill=C_GREEN)
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(idx + 1)
        run.font.size  = Pt(19)
        run.font.bold  = True
        run.font.color.rgb = C_WHITE
        run.font.name  = 'Calibri'

        # Item text — starts after badge, sz=26pt
        textbox(slide, i(2.70), i(item_y), i(10), i(0.64),
                item, size=Pt(26), color=C_TEXT_DARK)
    return slide


def build_section(prs, c, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COVER)

    # VIB exact: x=1.42" y=3.79" w=17.33" h=1.31" sz=64pt
    textbox(slide, i(1.42), i(3.79), i(17.33), i(1.31),
            c.get('title', ''), size=Pt(64), bold=True, color=C_WHITE)
    return slide


def build_content(prs, c, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_CONTENT)
    content_header(slide, c.get('title', ''))

    body_y = i(1.25)
    body_h = SLIDE_H - body_y - i(0.45)    # leave room for www.fss.com.vn footer
    img    = c.get('image')
    body   = c.get('body', '')
    has_img = bool(img) and c['layout'] == 'content'

    if has_img:
        textbox(slide, i(0.40), body_y, i(8.50), body_h,
                body, size=Pt(17), color=C_TEXT_DARK)
        placeholder(slide, i(9.20), body_y, SLIDE_W - i(9.60), body_h - i(0.20), img)
    else:
        textbox(slide, i(0.40), body_y, SLIDE_W - i(0.80), body_h,
                body, size=Pt(17), color=C_TEXT_DARK)
    return slide


def build_table(prs, c, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_CONTENT)
    content_header(slide, c.get('title', ''))

    header = c.get('header', [])
    rows   = c.get('rows',   [])
    if not header and not rows:
        return slide

    n_cols = max(len(header), max((len(r) for r in rows), default=1))
    n_rows = 1 + len(rows)
    tbl_y  = i(1.30)
    tbl_h  = SLIDE_H - tbl_y - i(0.45)
    tbl_w  = SLIDE_W - i(0.80)

    table = slide.shapes.add_table(n_rows, n_cols, i(0.40), tbl_y, tbl_w, tbl_h).table

    for ci, txt in enumerate(header[:n_cols]):
        cell = table.cell(0, ci)
        cell.text = txt
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(16); run.font.bold = True
        run.font.color.rgb = C_WHITE; run.font.name = 'Calibri'
        cell.fill.solid(); cell.fill.fore_color.rgb = C_GREEN

    for ri, row_data in enumerate(rows):
        bg = RGBColor(0xF4, 0xF4, 0xF4) if ri % 2 == 0 else C_WHITE
        for ci in range(n_cols):
            cell = table.cell(ri + 1, ci)
            cell.text = row_data[ci] if ci < len(row_data) else ''
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(15); run.font.name = 'Calibri'
            run.font.color.rgb = C_TEXT_DARK
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
    return slide


def build_closing(prs, c, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COVER)

    # VIB exact: x=1.28" y=4.30" w=17.23" h=1.01" sz=38pt (4800)
    textbox(slide, i(1.28), i(4.30), i(17.23), i(1.01),
            c.get('title', 'Q&A'), size=Pt(38), bold=True, color=C_WHITE)
    return slide


# ── Dispatcher ─────────────────────────────────────────────────────────────────

BUILDERS = {
    'cover':        build_cover,
    'toc':          build_toc,
    'section':      build_section,
    'content':      build_content,
    'content-text': build_content,
    'table':        build_table,
    'closing':      build_closing,
}


# ── Converter ──────────────────────────────────────────────────────────────────

def convert(md_text, output_path):
    meta, slides_raw = parse_md(md_text)
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    placeholders_used = []
    for raw in slides_raw:
        c = parse_content(raw)
        builder = BUILDERS.get(c['layout'])
        if not builder:
            print(f"WARNING: Unknown layout '{c['layout']}', skipping.")
            continue
        builder(prs, c, meta)
        if c.get('image'):
            placeholders_used.append(c['image'])

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"✓ Saved: {out}  ({len(prs.slides)} slides)")

    if placeholders_used:
        print(f"\nImage placeholders to replace ({len(placeholders_used)}):")
        for idx, desc in enumerate(placeholders_used, 1):
            print(f"  {idx}. {desc}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input',  required=True, help='Input .md file')
    parser.add_argument('--output', required=True, help='Output .pptx file')
    args = parser.parse_args()
    convert(Path(args.input).read_text(encoding='utf-8'), args.output)


if __name__ == '__main__':
    main()
